from __future__ import annotations

import csv
import json
import mimetypes
from dataclasses import dataclass
from pathlib import Path

import pymupdf
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


@dataclass(slots=True)
class ExtractedDocument:
    text: str
    title: str | None = None


class DocumentExtractor:
    TEXT_SUFFIXES = {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".xml",
        ".yaml",
        ".yml",
        ".log",
        ".py",
        ".js",
        ".ts",
        ".html",
        ".htm",
    }

    def extract(self, path: Path, mime_type: str | None = None) -> ExtractedDocument:
        suffix = path.suffix.lower()
        mime_type = mime_type or mimetypes.guess_type(path.name)[0]
        if suffix == ".pdf" or mime_type == "application/pdf":
            return self._pdf(path)
        if suffix == ".docx":
            return self._docx(path)
        if suffix in {".xlsx", ".xlsm"}:
            return self._xlsx(path)
        if suffix == ".pptx":
            return self._pptx(path)
        if suffix in self.TEXT_SUFFIXES or (mime_type and mime_type.startswith("text/")):
            return self._text(path)
        return ExtractedDocument(text="")

    @staticmethod
    def _pdf(path: Path) -> ExtractedDocument:
        document = pymupdf.open(path)
        try:
            metadata = document.metadata or {}
            pages = [page.get_text("text") for page in document]
            return ExtractedDocument(
                text="\n\n".join(page.strip() for page in pages if page.strip()),
                title=metadata.get("title") or None,
            )
        finally:
            document.close()

    @staticmethod
    def _docx(path: Path) -> ExtractedDocument:
        document = Document(path)
        lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                lines.append(" | ".join(cell.text.strip() for cell in row.cells))
        title = document.core_properties.title or None
        return ExtractedDocument(text="\n".join(lines), title=title)

    @staticmethod
    def _xlsx(path: Path) -> ExtractedDocument:
        workbook = load_workbook(path, read_only=True, data_only=True)
        lines: list[str] = []
        try:
            for worksheet in workbook.worksheets:
                lines.append(f"Sheet: {worksheet.title}")
                for row in worksheet.iter_rows(values_only=True):
                    values = [str(value) if value is not None else "" for value in row]
                    if any(values):
                        lines.append(" | ".join(values))
        finally:
            workbook.close()
        return ExtractedDocument(text="\n".join(lines), title=path.stem)

    @staticmethod
    def _pptx(path: Path) -> ExtractedDocument:
        presentation = Presentation(path)
        lines: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"Slide {index}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return ExtractedDocument(text="\n".join(lines), title=path.stem)

    @staticmethod
    def _text(path: Path) -> ExtractedDocument:
        raw = path.read_text(encoding="utf-8", errors="replace")
        if path.suffix.lower() == ".json":
            try:
                raw = json.dumps(json.loads(raw), ensure_ascii=False, indent=2)
            except json.JSONDecodeError:
                pass
        elif path.suffix.lower() == ".csv":
            try:
                rows = csv.reader(raw.splitlines())
                raw = "\n".join(" | ".join(row) for row in rows)
            except csv.Error:
                pass
        return ExtractedDocument(text=raw, title=path.stem)
